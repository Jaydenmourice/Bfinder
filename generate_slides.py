from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── App-exact colors ─────────────────────────────────────────────────────────
DARK_BG    = RGBColor(12,  16,  32)   # rgb(12,16,32)
PANEL_BG   = RGBColor(20,  28,  54)   # rgb(20,28,54)
CARD_BG    = RGBColor(26,  36,  68)   # rgb(26,36,68)
ACCENT     = RGBColor(0x4a, 0x9e, 0xff)  # #4a9eff
ACCENT_LT  = RGBColor(0x7a, 0xb8, 0xff)  # #7ab8ff  (lighter blue)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

LOGO = r"C:\Users\HP\Desktop\Bfinder\images\logo.png"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def box(slide, l, t, w, h, fill=PANEL_BG, border=None, radius=False):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def logo(slide, l, t, w, h):
    slide.shapes.add_picture(LOGO, Inches(l), Inches(t), Inches(w), Inches(h))


def accent_bar(slide, t=1.05, h=0.055):
    box(slide, 0, t, 13.33, h, fill=ACCENT)


def page_num(slide, n):
    txt(slide, str(n), 12.6, 7.1, 0.5, 0.3, size=10, color=LIGHT_GRAY,
        align=PP_ALIGN.RIGHT)


# ── Slide 1 · Title ──────────────────────────────────────────────────────────
s = add_slide(); bg(s)
# side accent stripe
box(s, 0, 0, 0.35, 7.5, fill=ACCENT)
box(s, 0.35, 0, 0.08, 7.5, fill=CARD_BG)
# content area
box(s, 0.6, 1.8, 12.5, 3.9, fill=PANEL_BG)
box(s, 0.6, 1.8, 12.5, 0.07, fill=ACCENT)   # top rule
box(s, 0.6, 5.63, 12.5, 0.07, fill=ACCENT)  # bottom rule
# logo
logo(s, 1.1, 2.0, 1.5, 1.5)
# title text
txt(s, "BFINDER", 2.9, 2.0, 9.8, 1.5, size=68, bold=True, color=ACCENT)
txt(s, "Web Security Vulnerability Scanner", 2.9, 3.5, 9.8, 0.75,
    size=24, color=WHITE)
txt(s, "Presenter: Jayden  ·  June 2026", 2.9, 4.3, 9.8, 0.5,
    size=15, color=LIGHT_GRAY, italic=True)
txt(s, "Supervisor Presentation", 2.9, 4.85, 9.8, 0.45,
    size=13, color=ACCENT_LT, italic=True)
page_num(s, 1)

# ── Slide 2 · Agenda ─────────────────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "AGENDA", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
items = [
    "01  Problem Statement",
    "02  What is Bfinder?",
    "03  Key Features",
    "04  System Architecture",
    "05  Live Demo Overview",
    "06  Technical Stack",
    "07  Future Roadmap",
    "08  Q & A",
]
for i, item in enumerate(items):
    cx = 0.5 if i < 4 else 7.0
    cy = 1.35 + (i % 4) * 1.45
    box(s, cx, cy, 5.8, 1.2, fill=PANEL_BG, border=ACCENT)
    box(s, cx, cy, 0.15, 1.2, fill=ACCENT)
    txt(s, item, cx + 0.35, cy + 0.35, 5.3, 0.6, size=16, color=WHITE)
page_num(s, 2)

# ── Slide 3 · Problem Statement ──────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Problem Statement", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
problems = [
    ("Security vulnerabilities slip to production",
     "Developers focus on features; security checks are often skipped or done too late."),
    ("Manual code review is slow & inconsistent",
     "Human reviewers miss subtle issues such as XSS, SQL injection, and exposed secrets."),
    ("No lightweight desktop tool for offline scanning",
     "Existing tools are cloud-dependent, expensive, or require DevOps-level setup."),
]
for i, (title, desc) in enumerate(problems):
    y = 1.25 + i * 1.9
    box(s, 0.5, y, 12.3, 1.65, fill=PANEL_BG)
    box(s, 0.5, y, 0.18, 1.65, fill=ACCENT)
    txt(s, title, 0.85, y + 0.1,  11.6, 0.55, size=17, bold=True, color=WHITE)
    txt(s, desc,  0.85, y + 0.7,  11.6, 0.8,  size=14, color=LIGHT_GRAY)
page_num(s, 3)

# ── Slide 4 · What is Bfinder? ───────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "What is Bfinder?", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)

box(s, 0.5, 1.25, 5.9, 5.75, fill=PANEL_BG)
txt(s, "Bfinder is a desktop application that automatically scans web project "
       "directories for common security vulnerabilities — giving developers "
       "instant, actionable results before deployment.",
    0.7, 1.4, 5.5, 2.2, size=14, color=WHITE)
txt(s, "Target Users", 0.7, 3.55, 5.5, 0.5, size=15, bold=True, color=ACCENT)
for i, u in enumerate(["Web Developers", "QA Engineers",
                        "Security Auditors", "CS Students"]):
    box(s, 0.75, 4.1 + i * 0.67, 0.28, 0.28, fill=ACCENT)
    txt(s, u, 1.15, 4.08 + i * 0.67, 4.9, 0.45, size=13, color=LIGHT_GRAY)

box(s, 7.1, 1.25, 5.9, 5.75, fill=PANEL_BG)
txt(s, "Goals", 7.3, 1.4, 5.5, 0.5, size=15, bold=True, color=ACCENT)
goals = [
    "Detect vulnerabilities early in development",
    "Work fully offline — no internet needed",
    "Simple UI that any developer can use",
    "Generate printable encrypted PDF reports",
    "Track scan history across 50 scans",
]
for i, g in enumerate(goals):
    box(s, 7.3, 2.05 + i * 1.0, 0.28, 0.28, fill=ACCENT)
    txt(s, g, 7.75, 2.02 + i * 1.0, 5.1, 0.6, size=13, color=WHITE)
page_num(s, 4)

# ── Slide 5 · Key Features ───────────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Key Features", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
features = [
    ("Vulnerability Scanner",  "Scans .html .js .ts .php .css .json .env .xml .sql for security bugs"),
    ("Login & Auth",           "Secure login with OS keyring-stored password and logout support"),
    ("Bug Explanations",       "Each detected issue comes with a description and fix suggestion"),
    ("PDF Reports",            "Generates encrypted PDF reports with full scan results"),
    ("Scan History",           "Stores last 50 scans in scan_history.json for later review"),
    ("Security Tips",          "Rotating security tips displayed on the home page at launch"),
]
positions = [(0.5,1.25),(4.6,1.25),(8.7,1.25),(0.5,4.1),(4.6,4.1),(8.7,4.1)]
for (cx, cy), (ft, fd) in zip(positions, features):
    box(s, cx, cy, 3.8, 2.6, fill=CARD_BG, border=ACCENT)
    box(s, cx, cy, 3.8, 0.07, fill=ACCENT)
    txt(s, ft, cx+0.15, cy+0.18, 3.5, 0.6, size=14, bold=True, color=ACCENT)
    txt(s, fd, cx+0.15, cy+0.82, 3.5, 1.6, size=12, color=LIGHT_GRAY)
page_num(s, 5)

# ── Slide 6 · System Architecture ───────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "System Architecture", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
layers = [
    ("UI Layer  (PyQt5)",      "LoginDialog  ·  BfinderWindow  ·  QStackedWidget Pages",   ACCENT),
    ("Logic Layer",            "bug_exp.py — scan engine  ·  sec_tip.py — tips engine",     ACCENT_LT),
    ("Data Layer",             "scan_history.json  ·  contact_info.json  ·  PDF Reports",   RGBColor(0x5b, 0xb8, 0xff)),
    ("OS / Security Layer",    "keyring (password storage)  ·  xhtml2pdf + PyPDF2 (encryption)", RGBColor(0x30, 0x70, 0xcc)),
]
for i, (title, detail, col) in enumerate(layers):
    y = 1.3 + i * 1.45
    box(s, 0.5, y, 12.3, 1.25, fill=PANEL_BG)
    box(s, 0.5, y, 3.0,  1.25, fill=col)
    txt(s, title,  0.65, y+0.38, 2.75, 0.6,  size=14, bold=True, color=DARK_BG)
    txt(s, detail, 3.8,  y+0.38, 8.8,  0.6,  size=13, color=LIGHT_GRAY)
page_num(s, 6)

# ── Slide 7 · Supported File Types ──────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Supported File Types", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
types = [
    (".html",    "XSS, inline scripts, unsafe forms"),
    (".js / .ts","eval(), dangerous sinks, exposed keys"),
    (".php",     "SQL injection, exec(), file inclusion"),
    (".env",     "Exposed secrets, API keys, passwords"),
    (".sql",     "Unsafe queries, DROP/TRUNCATE risks"),
    (".json",    "Hardcoded credentials, insecure configs"),
    (".css",     "CSS injection, expression() exploits"),
    (".xml",     "XXE injection, external entity refs"),
]
for i, (ext, desc) in enumerate(types):
    cx = 0.5 if i < 4 else 7.0
    cy = 1.3 + (i % 4) * 1.52
    box(s, cx, cy, 5.8, 1.3, fill=CARD_BG, border=ACCENT)
    txt(s, ext,  cx+0.2, cy+0.12, 1.5, 0.5, size=15, bold=True, color=ACCENT)
    txt(s, desc, cx+1.9, cy+0.13, 3.7, 0.8, size=13, color=WHITE)
page_num(s, 7)

# ── Slide 8 · Technical Stack ────────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Technical Stack", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
stack = [
    ("Python 3",    "Core language — cross-platform, rich ecosystem"),
    ("PyQt5",       "Desktop UI framework — windows, widgets, event loop"),
    ("keyring",     "OS-native secure credential storage"),
    ("xhtml2pdf",   "HTML-to-PDF conversion for report generation"),
    ("PyPDF2",      "PDF encryption and manipulation"),
    ("JSON",        "Lightweight local storage for history and settings"),
]
for i, (tech, desc) in enumerate(stack):
    y = 1.3 + i * 1.0
    box(s, 0.5,  y, 12.3, 0.82, fill=PANEL_BG)
    box(s, 0.5,  y, 2.6,  0.82, fill=ACCENT)
    txt(s, tech, 0.65, y+0.18, 2.3, 0.5, size=14, bold=True, color=DARK_BG)
    txt(s, desc, 3.35, y+0.18, 9.2, 0.5, size=13, color=WHITE)
page_num(s, 8)

# ── Slide 9 · Demo Overview ──────────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Demo Overview", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
steps = [
    ("1", "Launch Bfinder",        "Enter password to unlock the application"),
    ("2", "Browse Project Folder", "Select a web project directory to scan"),
    ("3", "Run Scan",              "Click Scan — Bfinder analyses all supported files"),
    ("4", "Review Results",        "View detected bugs with explanations on Home page"),
    ("5", "Generate Report",       "Click Print Report to produce an encrypted PDF"),
    ("6", "Check History",         "Navigate to History to review previous scans"),
]
for i, (num, title, desc) in enumerate(steps):
    cx = 0.5 if i < 3 else 7.0
    cy = 1.25 + (i % 3) * 2.05
    box(s, cx, cy, 5.8, 1.8, fill=CARD_BG, border=ACCENT)
    box(s, cx, cy, 1.1, 1.8, fill=ACCENT)
    txt(s, num,   cx+0.25, cy+0.6,  0.7, 0.7, size=28, bold=True, color=DARK_BG,
        align=PP_ALIGN.CENTER)
    txt(s, title, cx+1.3,  cy+0.12, 4.3, 0.6, size=15, bold=True, color=WHITE)
    txt(s, desc,  cx+1.3,  cy+0.75, 4.3, 0.85, size=13, color=LIGHT_GRAY)
page_num(s, 9)

# ── Slide 10 · Future Roadmap ────────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s)
logo(s, 0.35, 0.15, 0.75, 0.75)
txt(s, "Future Roadmap", 1.25, 0.18, 11, 0.7, size=32, bold=True, color=ACCENT)
phases = [
    ("Phase 1\nShort Term",  ["More scan rules & patterns", "Support additional file types",  "Better bug explanation detail"]),
    ("Phase 2\nMid Term",    ["CI/CD pipeline integration",  "Command-line interface (CLI)",   "Team-shared scan history"]),
    ("Phase 3\nLong Term",   ["Web / cloud-based version",   "AI-powered fix suggestions",     "Real-time scan as-you-type"]),
]
phase_colors = [ACCENT, ACCENT_LT, RGBColor(0x30, 0x70, 0xcc)]
for i, (phase, items) in enumerate(phases):
    x = 0.5 + i * 4.3
    box(s, x, 1.3, 4.1, 5.8, fill=PANEL_BG)
    box(s, x, 1.3, 4.1, 1.15, fill=phase_colors[i])
    txt(s, phase, x+0.1, 1.35, 3.9, 1.05, size=16, bold=True,
        color=DARK_BG, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        box(s, x+0.3, 2.7+j*1.35, 0.3, 0.3, fill=phase_colors[i])
        txt(s, item, x+0.8, 2.65+j*1.35, 3.1, 0.65, size=13, color=WHITE)
page_num(s, 10)

# ── Slide 11 · Thank You / Q&A ───────────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.35, 7.5, fill=ACCENT)
box(s, 0.35, 0, 0.08, 7.5, fill=CARD_BG)
box(s, 0.6, 1.7, 12.4, 4.1, fill=PANEL_BG)
box(s, 0.6, 1.7, 12.4, 0.07, fill=ACCENT)
box(s, 0.6, 5.73, 12.4, 0.07, fill=ACCENT)
logo(s, 1.2, 2.0, 1.6, 1.6)
txt(s, "Thank You", 3.1, 1.95, 9.7, 1.4, size=58, bold=True, color=ACCENT)
txt(s, "Questions & Discussion", 3.1, 3.4, 9.7, 0.75, size=22, color=WHITE)
txt(s, "waltermkingilima96@gmail.com", 3.1, 4.25, 9.7, 0.5, size=15, color=ACCENT_LT)
txt(s, "Bfinder  ·  Web Security Scanner  ·  June 2026", 0, 6.3, 13.33, 0.4,
    size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
page_num(s, 11)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\HP\Desktop\Bfinder\Bfinder_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
